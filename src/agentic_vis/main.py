"""Main module for agentic_vis.

Generates PowerPoint presentations with agent progress visualization.
"""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pathlib import Path

import win32com

from agentic_vis.openers import full_control_open_pptx
from src.agentic_vis.harvey import generate_harvey_balls
from src.agentic_vis.draw import center_tool_nodes,center_l1_nodes,create_drawing_slide, connect_circles, draw_circles_center, resolve_circle_overlaps, write_surround_circles
from src.agentic_vis.excel_utils import attach_harvey_balls, parse_excel_l12a, parse_excel_a2t, parse_excel_components

from src.agentic_vis.excel_utils import attach_harvey_balls

PROJ_ROOT = 'C:\\Users\\John_Harney\\github\\graph_tutorials'
XLSX_FILE_PATH = "C:\\Users\\John_Harney\\github\\graph_tutorials\\hierarchy.xlsx"

AGENTIC_VIS_TITLE = "Agentic Visualization Dashboard1"
DEST_FILE_DIR = PROJ_ROOT
DEST_FILE_NAME = 'agentic_plan.pptx'
SOURCE_FILE_DIR = PROJ_ROOT + '\\private_templates\\'
SOURCE_FILE_NAME = "slide_template_blank.pptx"
SLIDE_TO_COPY = 3


# `create_drawing_slide` is implemented in src/agentic_vis/draw.py



def create_presentation(prs,filename: str = "agentic_plan.pptx"):

    items = ['a','b','c','aa','bb','cc','a']
    sizes = [1, 2, 1, 1, 1, 1,1]
    slide = create_drawing_slide(prs, items, sizes)
    harvey_ball_dict = {
        'a' : 2,
        'aa' : 2,
        'b' : 1,
        'bb' : 3,
        'c' : 4
    }
    connect_circles(slide,'a','b')
    #generate_harvey_balls(slide, 

    #    ['a', 'c']
    #)
    generate_harvey_balls(
        slide,
        harvey_ball_dict,
    )
    prs.save(filename)


from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import math

from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import math
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import math

def draw_circles_around_source(slide, source: str, dest: list[str]) -> None:
    """
    Draw 6 circles of size 0.25 around the source circle.
    """
    pass


def personal_main2(
        pptx_output_file_path: str = "C:\\Users\\jfhar\\OneDrive\\Desktop\\github\\graph_tutorials\\agentic_plan.pptx",
        xlsx_input_file_path: str = "C:\\Users\\jfhar\\OneDrive\\Desktop\\github\\graph_tutorials\\hierarchy.xlsx"
    ) -> None:
    l12a_pairs = parse_excel_l12a(xlsx_input_file_path)
    print ('l12a_pairs:')
    print (l12a_pairs)
    a2t_pairs = parse_excel_a2t(xlsx_input_file_path)
    print ('a2t_pairs:')
    print (a2t_pairs)
    components_response = parse_excel_components(xlsx_input_file_path)
    
    print ('components_response:')
    print (components_response)
    items = components_response[0]
    sizes = components_response[1]
    print (items)
    print (sizes)
    
    prs = Presentation()
    slide = create_drawing_slide(prs, items, sizes)
    

    
    items = components_response[0]
    sizes = components_response[1]
    print (items)
    print (sizes)
    center_l1_nodes(slide)
    center_tool_nodes(slide)

    resolve_circle_overlaps(slide)

    write_surround_circles(slide, "Personalization", ["A1","A2","A3","A4","A5"])
    write_surround_circles(slide, "Knowledge", ["A6","A7","A8","A9","A10"])


    harvey_ball_dict = attach_harvey_balls(xlsx_input_file_path)
    generate_harvey_balls(
        slide,
        harvey_ball_dict,
    )

    prs.save(pptx_output_file_path)
    print (f"saved {pptx_output_file_path}")


def main() -> None:
    """Entry point for agentic_vis."""
    
    
    pptx_output_file_path = "C:\\Users\\jfhar\\OneDrive\\Desktop\\github\\graph_tutorials\\agentic_plan.pptx"
    xlsx_input_file_path: str = "C:\\Users\\jfhar\\OneDrive\\Desktop\\github\\graph_tutorials\\hierarchy.xlsx"
    
    full_control_open_pptx(pptx_output_file_path)
    personal_main2(
        pptx_output_file_path,
        xlsx_input_file_path,
    )
    

    #full_control_open_pptx(pptx_output_file_path, False)

    # do something here to create a new slide deck

    full_control_open_pptx(pptx_output_file_path)

if __name__ == "__main__":
    main()


