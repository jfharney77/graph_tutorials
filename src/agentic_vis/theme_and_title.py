
from pptx.util import Inches, Pt
from pptx import Presentation
from pptx.enum.text import PP_ALIGN

def use_previous_slide_deck(file_path: str) -> Presentation:
    return Presentation(file_path)

def apply_theme(presentation, slide_num):
    slide = presentation.slides[slide_num - 1]  # python-pptx uses 0-based indexing
    return slide.slide_layout

def create_slide_from_theme(prs, layout):
    return prs.slides.add_slide(layout)

def add_title_to_individual_slide(s, title_text):
    if hasattr(s, "shapes") and hasattr(s.shapes, "title"):
        s.shapes.title.text = title_text
    else:
        left = Inches(2.5)
        top = Inches(0.5)
        width = Inches(5)
        height = Inches(1)
        txBox = s.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = title_text
        tf.paragraphs[0].font.size = Pt(24)



def add_title_to_slide(s, title_name: str) -> None:
    left = Inches(2.5)
    top = Inches(0.5)
    width = Inches(5)
    height = Inches(1)
    txBox = s.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = title_name
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(24)

