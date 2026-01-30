"""Harvey ball generation for PowerPoint presentations.

Supports fill levels 0-4 where:
 - 0 = empty
 - 1 = 1/4 filled (bottom-right)
 - 2 = 1/2 filled (bottom half)
 - 3 = 3/4 filled (all but top-left)
 - 4 = fully filled
"""
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor


def _add_rect(slide, left, top, width, height, color=RGBColor(0, 0, 0)):
    """Helper to add a filled rectangle without visible border."""
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(left), int(top), int(width), int(height))
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    return rect


def generate_harvey_balls(slide, items: dict[str, int]) -> None:
    """Draw small indicator circles (Harvey balls) to the top-left of specified circles.

    Args:
        slide: The slide object returned from create_drawing_slide.
        items: Dictionary where keys are circle names and values are integers (0-4).
    """
    if not hasattr(slide, 'circle_info'):
        print("Error: No circle information found on slide. Use slide from create_drawing_slide().")
        return

    circle_info = slide.circle_info
    harvey_ball_size = Inches(0.125)
    offset_pixels = 10

    for item_name, fill_value in items.items():
        if item_name not in circle_info:
            print(f"Warning: Circle '{item_name}' not found on slide, skipping Harvey ball.")
            continue

        # Get circle info
        left, top, size, cx, cy = circle_info[item_name]

        # Position Harvey ball to the top-left of the circle
        harvey_left = left - harvey_ball_size - offset_pixels
        harvey_top = top - harvey_ball_size - offset_pixels

        # Create outer Harvey ball (always present)
        harvey_ball = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            int(harvey_left),
            int(harvey_top),
            harvey_ball_size,
            harvey_ball_size,
        )

        # Decide fill
        if fill_value == 0:
            # empty
            harvey_ball.fill.background()
        elif fill_value == 4:
            # full
            harvey_ball.fill.solid()
            harvey_ball.fill.fore_color.rgb = RGBColor(0, 0, 0)
        elif fill_value == 2:
            # half (bottom half)
            harvey_ball.fill.background()
            half_h = harvey_ball_size / 2
            _add_rect(slide, harvey_left, harvey_top + half_h, harvey_ball_size, half_h)
        elif fill_value == 1:
            # quarter (bottom-right)
            harvey_ball.fill.background()
            half_w = harvey_ball_size / 2
            half_h = harvey_ball_size / 2
            _add_rect(slide, harvey_left + half_w, harvey_top + half_h, half_w, half_h)
        elif fill_value == 3:
            # three quarters: bottom half + top-right quarter
            harvey_ball.fill.background()
            half_w = harvey_ball_size / 2
            half_h = harvey_ball_size / 2
            _add_rect(slide, harvey_left, harvey_top + half_h, harvey_ball_size, half_h)
            _add_rect(slide, harvey_left + half_w, harvey_top, half_w, half_h)
        else:
            harvey_ball.fill.background()

        # Border color
        harvey_ball.line.color.rgb = RGBColor(184, 134, 11)


"""Harvey ball generation for PowerPoint presentations."""
'''
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor


def generate_harvey_balls(slide, items: dict[str, int]) -> None:
    """Draw small indicator circles (Harvey balls) to the top-left of specified circles.

    Args:
        slide: The slide object returned from create_drawing_slide.
        items: Dictionary where keys are circle names and values are integers.
               If value is 0: unfilled circle
               If value is 2: half-filled circle (black bottom half)
               If value is 4: fully filled circle (black)
    """
    if not hasattr(slide, 'circle_info'):
        print("Error: No circle information found on slide. Use slide from create_drawing_slide().")
        return

    circle_info = slide.circle_info
    harvey_ball_size = Inches(0.125)
    offset_pixels = 10

    for item_name, fill_value in items.items():
        if item_name not in circle_info:
            print(f"Warning: Circle '{item_name}' not found on slide, skipping Harvey ball.")
            continue

        # Get circle info
        left, top, size, cx, cy = circle_info[item_name]
        radius = size / 2

        # Position Harvey ball to the top-left of the circle
        # 10 pixels to the left of the left edge and 10 pixels above the top edge
        harvey_left = left - harvey_ball_size - offset_pixels
        harvey_top = top - harvey_ball_size - offset_pixels

        # Create Harvey ball (outer circle)
        harvey_ball = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            int(harvey_left),
            int(harvey_top),
            harvey_ball_size,
            harvey_ball_size
        )

        # Format Harvey ball based on fill_value
        if fill_value == 0:
            # No fill
            harvey_ball.fill.background()
        elif fill_value == 4:
            # Full fill with black
            harvey_ball.fill.solid()
            harvey_ball.fill.fore_color.rgb = RGBColor(0, 0, 0)  # Black fill
        elif fill_value == 2:
            # Half fill - create empty circle, then overlay a rectangle for bottom half
            harvey_ball.fill.background()
            # Add a rectangle covering the bottom half
            half_height = harvey_ball_size / 2
            rect = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                int(harvey_left),
                int(harvey_top + half_height),
                int(harvey_ball_size),
                int(half_height)
            )
            rect.fill.solid()
            rect.fill.fore_color.rgb = RGBColor(0, 0, 0)  # Black fill
            rect.line.color.rgb = RGBColor(0, 0, 0)  # No border or matching color
        else:
            # Default to no fill for unknown values
            harvey_ball.fill.background()

        # Border color
        harvey_ball.line.color.rgb = RGBColor(184, 134, 11)  # Dark goldenrod border
'''