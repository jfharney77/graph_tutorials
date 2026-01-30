from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import math

from src.agentic_vis.excel_utils import TOOL_SIZE, L1_SIZE, AGENT_SIZE
# (old version from commit e6773f3)

def write_surround_circles(slide, source: str, dest: list) -> None:
    """Position or create circles around a source circle based on dest list.

    For each item name in dest, if a circle with that name exists, move it to
    1.5 inches from the source center. If it doesn't exist, create a new circle
    of size AGENT_SIZE at that position.

    Args:
        slide: The slide object with circle_info attribute.
        source: The name of the source circle to surround.
        dest: A list of circle names to position around the source.
    """
    if not hasattr(slide, 'circle_info') or not isinstance(slide.circle_info, dict):
        print("Error: No circle information found on slide.")
        return

    circle_info = slide.circle_info

    if source not in circle_info:
        print(f"Error: Circle '{source}' not found on slide.")
        return

    # Get source circle info
    left_src, top_src, size_src, cx_src, cy_src = circle_info[source]

    print ('cx_src ' + str(cx_src) + ' cy_src ' + str(cy_src))

    radius_distance = Inches(1.5)  # 1.5 inches from source center

    # Process each destination circle
    num_circles = len(dest)
    if num_circles == 0:
        return  # No circles to process

    for i, dest_name in enumerate(dest):
        angle = (2 * math.pi * i) / num_circles
        # Calculate position at 1.5 inches from source center
        x_offset = radius_distance * math.cos(angle)
        y_offset = radius_distance * math.sin(angle)

        # Calculate center position for the circle
        new_cx = cx_src + x_offset
        new_cy = cy_src + y_offset
        print ('new_cx ' + str(new_cx) + ' new_cy ' + str(new_cy))

        # Check if circle with dest_name already exists
        if dest_name in circle_info:
            # Circle exists - move it
            _, _, existing_size, _, _ = circle_info[dest_name]
            new_left = new_cx - existing_size / 2
            new_top = new_cy - existing_size / 2

            # Update circle_info
            circle_info[dest_name] = (new_left, new_top, existing_size, new_cx, new_cy)

            # Find and move the shape
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame.text == dest_name:
                    shape.left = int(new_left)
                    shape.top = int(new_top)
                    break
        else:
            # Circle doesn't exist - create new one with AGENT_SIZE
            circle_size = Inches(AGENT_SIZE)
            new_left = new_cx - circle_size / 2
            new_top = new_cy - circle_size / 2

            # Add the circle shape
            shape = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                int(new_left),
                int(new_top),
                circle_size,
                circle_size
            )

            # Format the circle
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(100, 149, 237)  # Cornflower blue
            shape.line.color.rgb = RGBColor(25, 25, 112)  # Midnight blue

            # Add text to circle
            text_frame = shape.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            p.text = dest_name
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)  # White text
            p.alignment = 1  # Center alignment

            # Store in circle_info
            circle_info[dest_name] = (new_left, new_top, circle_size, new_cx, new_cy)

