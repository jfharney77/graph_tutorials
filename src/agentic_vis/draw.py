from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import math

from src.agentic_vis.excel_utils import TOOL_SIZE, L1_SIZE, AGENT_SIZE


def connect_circles(slide, circ1_name: str, circ2_name: str) -> None:
    """Draw a line connecting two circles by their names at their circumferences.

    Args:
        slide: The slide object returned from create_drawing_slide.
        circ1_name: The text name of the first circle.
        circ2_name: The text name of the second circle.
    """
    if not hasattr(slide, 'circle_info'):
        print("Error: No circle information found on slide. Use slide from create_drawing_slide().")
        return

    circle_info = slide.circle_info

    #print ('exiting...')
    #import sys
    #sys.exit()

    if circ1_name not in circle_info or circ2_name not in circle_info:
        print(f"Error: Circle '{circ1_name}' or '{circ2_name}' not found on slide.")
        return

    # Get center coordinates and sizes of both circles
    left1, top1, size1, cx1, cy1 = circle_info[circ1_name]
    left2, top2, size2, cx2, cy2 = circle_info[circ2_name]

    # Calculate radii (half of size)
    radius1 = size1 / 2
    radius2 = size2 / 2

    # Calculate angle from circle1 center to circle2 center
    dx = cx2 - cx1
    dy = cy2 - cy1
    distance = math.sqrt(dx**2 + dy**2)

    if distance == 0:
        print("Warning: Circles have the same center, cannot draw a line.")
        return

    # Unit vector pointing from circle1 to circle2
    unit_dx = dx / distance
    unit_dy = dy / distance

    # Start point: on circumference of circle1 in direction of circle2
    start_x = cx1 + unit_dx * radius1
    start_y = cy1 + unit_dy * radius1

    # End point: on circumference of circle2 in direction of circle1
    end_x = cx2 - unit_dx * radius2
    end_y = cy2 - unit_dy * radius2

    # Draw a straight connector from circumference to circumference
    connector = slide.shapes.add_connector(1, int(start_x), int(start_y), int(end_x), int(end_y))
    connector.line.color.rgb = RGBColor(128, 128, 128)  # Gray line
    connector.line.width = Pt(2)  # 2 point width

def create_drawing_slide_from_slide(prs,slide, items: list[str] | None = None, sizes: list[float] | None = None):
    # Don't draw anything if items is None or empty list
    if not items or len(items) == 0:
        return slide

    # If a sizes list is provided but empty, don't draw anything per spec
    if sizes is not None and len(sizes) == 0:
        return slide

    # Deduplicate items while preserving order: if an item appears multiple
    # times in the input list, only draw the first occurrence.
    seen = set()
    unique_items = []
    for it in items:
        if it not in seen:
            seen.add(it)
            unique_items.append(it)

    # Get slide dimensions
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Set default sizes if not provided. Map sizes to the first occurrence
    # of each unique item so duplicates use the first size provided.
    if sizes is None:
        sizes = [Inches(1.5)] * len(unique_items)
    else:
        orig_sizes = [Inches(s) if isinstance(s, (int, float)) else s for s in sizes]
        sizes = []
        for it in unique_items:
            try:
                idx = items.index(it)
                sizes.append(orig_sizes[idx] if idx < len(orig_sizes) else Inches(1.5))
            except ValueError:
                sizes.append(Inches(1.5))

    num_items = len(unique_items)

    # Calculate spacing
    if num_items == 1:
        # Center single circle
        circle_size = sizes[0]
        left = (slide_width - circle_size) / 2
        top = (slide_height - circle_size) / 2
        positions = [(left, top)]
    else:
        # Arrange in a circle pattern around the center
        import math
        center_x = slide_width / 2
        center_y = slide_height / 2
        radius = Inches(2)

        positions = []
        for i in range(num_items):
            angle = (2 * math.pi * i) / num_items
            circle_size = sizes[i]
            x = center_x + radius * math.cos(angle) - circle_size / 2
            y = center_y + radius * math.sin(angle) - circle_size / 2
            positions.append((x, y))

    # Store circle info: name -> (left, top, size, center_x, center_y)
    circle_info = {}

    # Create circles with text
    for i, (item, (left, top)) in enumerate(zip(unique_items, positions)):
        circle_size = sizes[i]
        center_x = left + circle_size / 2
        center_y = top + circle_size / 2
        circle_info[item] = (left, top, circle_size, center_x, center_y)

        # Add circle shape
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            int(left),
            int(top),
            circle_size,
            circle_size
        )

        # Format circle
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGBColor(100, 149, 237)  # Cornflower blue
        circle.line.color.rgb = RGBColor(25, 25, 112)  # Midnight blue

        # Add text to circle
        text_frame = circle.text_frame
        text_frame.clear()  # Clear default text
        p = text_frame.paragraphs[0]
        p.text = item
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)  # White text
        p.alignment = 1  # Center alignment

    # Store circle info on slide for later use
    slide.circle_info = circle_info
    return slide


def connect_circles_batch(slide, pairs: list[tuple]) -> None:
    """Connect multiple pairs of circles in batch.

    For each tuple pair in the list, calls connect_circles with the pair's
    first element as circ1_name and second element as circ2_name.

    Args:
        slide: The slide object (typically returned from create_drawing_slide).
        pairs: A list of tuples where each tuple is (circ1_name, circ2_name).
    """
    for pair in pairs:
        if len(pair) >= 2:
            print ('connecting ' + str(pair[0] + ' ' + str(pair[1])))
            connect_circles(slide, pair[0], pair[1])


def create_drawing_slide(prs, items: list[str] | None = None, sizes: list[float] | None = None):
    """Create a slide with circles containing text for each item in the list.

    Args:
        prs: The Presentation object to add the slide to.
        items: Optional list of strings to display in circles. If None, creates empty slide.
        sizes: Optional list of circle sizes (in inches) matching the length of items.
               If None, defaults to 1.5 inches for all circles.

    Returns:
        The slide object with `circle_info` attribute mapping names to positions/sizes,
        or `None` if no items were provided.
    """
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_slide_layout)

    # Don't draw anything if items is None or empty list
    if not items or len(items) == 0:
        return slide

    # If a sizes list is provided but empty, don't draw anything per spec
    if sizes is not None and len(sizes) == 0:
        return slide

    # Deduplicate items while preserving order: if an item appears multiple
    # times in the input list, only draw the first occurrence.
    seen = set()
    unique_items = []
    for it in items:
        if it not in seen:
            seen.add(it)
            unique_items.append(it)

    # Get slide dimensions
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Set default sizes if not provided. Map sizes to the first occurrence
    # of each unique item so duplicates use the first size provided.
    if sizes is None:
        sizes = [Inches(1.5)] * len(unique_items)
    else:
        orig_sizes = [Inches(s) if isinstance(s, (int, float)) else s for s in sizes]
        sizes = []
        for it in unique_items:
            try:
                idx = items.index(it)
                sizes.append(orig_sizes[idx] if idx < len(orig_sizes) else Inches(1.5))
            except ValueError:
                sizes.append(Inches(1.5))

    num_items = len(unique_items)

    # Calculate spacing
    if num_items == 1:
        # Center single circle
        circle_size = sizes[0]
        left = (slide_width - circle_size) / 2
        top = (slide_height - circle_size) / 2
        positions = [(left, top)]
    else:
        # Arrange in a circle pattern around the center
        import math
        center_x = slide_width / 2
        center_y = slide_height / 2
        radius = Inches(2)

        positions = []
        for i in range(num_items):
            angle = (2 * math.pi * i) / num_items
            circle_size = sizes[i]
            x = center_x + radius * math.cos(angle) - circle_size / 2
            y = center_y + radius * math.sin(angle) - circle_size / 2
            positions.append((x, y))

    # Store circle info: name -> (left, top, size, center_x, center_y)
    circle_info = {}

    # Create circles with text
    for i, (item, (left, top)) in enumerate(zip(unique_items, positions)):
        circle_size = sizes[i]
        center_x = left + circle_size / 2
        center_y = top + circle_size / 2
        circle_info[item] = (left, top, circle_size, center_x, center_y)

        # Add circle shape
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            int(left),
            int(top),
            circle_size,
            circle_size
        )

        # Format circle
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGBColor(100, 149, 237)  # Cornflower blue
        circle.line.color.rgb = RGBColor(25, 25, 112)  # Midnight blue

        # Add text to circle
        text_frame = circle.text_frame
        text_frame.clear()  # Clear default text
        p = text_frame.paragraphs[0]
        p.text = item
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)  # White text
        p.alignment = 1  # Center alignment

    # Store circle info on slide for later use
    slide.circle_info = circle_info
    return slide


def add_circle_at_coords(slide, x: float, y: float) -> None:
    """Add a circle of size 0.5 inches at the specified coordinates on the slide.

    Args:
        slide: The slide object to add the circle to.
        x: The x coordinate (left position) in EMUs or as a numeric value.
        y: The y coordinate (top position) in EMUs or as a numeric value.
    """
    circle_size = Inches(0.5)

    # Add circle shape
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        int(x),
        int(y),
        circle_size,
        circle_size
    )

    # Format circle with default styling
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(100, 149, 237)  # Cornflower blue
    circle.line.color.rgb = RGBColor(25, 25, 112)  # Midnight blue


def draw_circles_center(slide, circle_num: int) -> None:
    """Draw `circle_num` circles of size 0.5 in a balanced layout centered on slide.

    Args:
        slide: The slide object to draw on.
        circle_num: Number of circles to draw (int).
    """
    if circle_num <= 0:
        return

    circle_size = Inches(0.5)

    # Determine slide dimensions
    prs = None
    try:
        prs = slide.part.package.presentation if hasattr(slide.part, 'package') else None
    except Exception:
        prs = None

    if prs is None:
        # fallback defaults (10in x 7.5in in EMUs)
        slide_width = Inches(10)
        slide_height = Inches(7.5)
    else:
        slide_width = prs.slide_width
        slide_height = prs.slide_height

    # margins
    margin_x = Inches(0.75)
    margin_y = Inches(0.75)

    avail_w = slide_width - 2 * margin_x
    avail_h = slide_height - 2 * margin_y

    # Compute grid size (rows x cols) as square as possible
    import math as _math
    cols = int(_math.ceil(_math.sqrt(circle_num)))
    rows = int(_math.ceil(circle_num / cols))

    # Compute spacing between circles
    if cols > 1:
        spacing_x = max((avail_w - cols * circle_size) / (cols - 1), Inches(0.1))
    else:
        spacing_x = 0
    if rows > 1:
        spacing_y = max((avail_h - rows * circle_size) / (rows - 1), Inches(0.1))
    else:
        spacing_y = 0

    # Starting top-left to center the grid
    total_w = cols * circle_size + (cols - 1) * spacing_x
    total_h = rows * circle_size + (rows - 1) * spacing_y

    start_x = (slide_width - total_w) / 2
    start_y = (slide_height - total_h) / 2

    # Ensure circle_info exists
    if not hasattr(slide, 'circle_info') or not isinstance(slide.circle_info, dict):
        print (type(slide))
        slide.circle_info = {}

    idx = 0
    for r in range(rows):
        for c in range(cols):
            if idx >= circle_num:
                break
            x = int(start_x + c * (circle_size + spacing_x))
            y = int(start_y + r * (circle_size + spacing_y))

            shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, circle_size, circle_size)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(100, 149, 237)
            shape.line.color.rgb = RGBColor(25, 25, 112)

            name = f"draw_center_{idx+1}"
            # store center coords in circle_info
            cx = x + circle_size / 2
            cy = y + circle_size / 2
            slide.circle_info[name] = (x, y, circle_size, cx, cy)

            idx += 1
        if idx >= circle_num:
            break


def resolve_circle_overlaps(slide, spacing: int = 20) -> None:
    """Check for overlapping circles on the slide and move them apart.

    This function examines all circles stored in slide.circle_info, detects
    overlaps based on distance between centers, and adjusts positions to
    eliminate overlaps with minimum spacing between outer circumferences.

    Args:
        slide: The slide object with circle_info attribute.
        spacing: Minimum pixel spacing between outer circumferences (default 20).
    """
    if not hasattr(slide, 'circle_info') or not isinstance(slide.circle_info, dict):
        print("Error: No circle information found on slide.")
        return

    circle_info = slide.circle_info
    if len(circle_info) < 2:
        return  # No overlaps possible with 0 or 1 circle

    # Build list of circle names
    circle_names = list(circle_info.keys())

    # Iterate multiple times to resolve cascading overlaps
    max_iterations = 10
    for iteration in range(max_iterations):
        overlap_found = False

        # Check all pairs
        for i in range(len(circle_names)):
            for j in range(i + 1, len(circle_names)):
                name1 = circle_names[i]
                name2 = circle_names[j]

                left1, top1, size1, cx1, cy1 = circle_info[name1]
                left2, top2, size2, cx2, cy2 = circle_info[name2]

                radius1 = size1 / 2
                radius2 = size2 / 2

                # Calculate distance between centers
                dx = cx2 - cx1
                dy = cy2 - cy1
                distance = math.sqrt(dx**2 + dy**2)

                # Check for overlap (include spacing between circumferences)
                min_distance = radius1 + radius2 + spacing
                if distance < min_distance and distance > 0:
                    overlap_found = True

                    # Calculate overlap amount
                    overlap = min_distance - distance

                    # Unit vector from circle1 to circle2
                    unit_dx = dx / distance
                    unit_dy = dy / distance

                    # Move each circle half the overlap distance apart
                    move_amount = (overlap / 2) * 1.1  # Add 10% padding

                    # Update positions for circle1 (move away from circle2)
                    new_cx1 = cx1 - unit_dx * move_amount
                    new_cy1 = cy1 - unit_dy * move_amount
                    new_left1 = new_cx1 - radius1
                    new_top1 = new_cy1 - radius1

                    # Update positions for circle2 (move away from circle1)
                    new_cx2 = cx2 + unit_dx * move_amount
                    new_cy2 = cy2 + unit_dy * move_amount
                    new_left2 = new_cx2 - radius2
                    new_top2 = new_cy2 - radius2

                    # Update circle_info
                    circle_info[name1] = (new_left1, new_top1, size1, new_cx1, new_cy1)
                    circle_info[name2] = (new_left2, new_top2, size2, new_cx2, new_cy2)

                    # Find and update the shapes
                    for shape in slide.shapes:
                        if hasattr(shape, 'text_frame') and shape.text_frame.text == name1:
                            shape.left = int(new_left1)
                            shape.top = int(new_top1)
                        elif hasattr(shape, 'text_frame') and shape.text_frame.text == name2:
                            shape.left = int(new_left2)
                            shape.top = int(new_top2)

        # If no overlaps found, we're done
        if not overlap_found:
            break


def center_l1_nodes(slide) -> None:
    """Find all circles of size 1.5 inches and arrange them in a row at slide center.

    This function identifies circles with size 1.5 inches in circle_info and
    repositions them horizontally centered on the slide.

    Args:
        slide: The slide object with circle_info attribute.
    """
    if not hasattr(slide, 'circle_info') or not isinstance(slide.circle_info, dict):
        print("Error: No circle information found on slide.")
        return

    circle_info = slide.circle_info

    # Find all circles with size 1.5 inches
    target_size = Inches(L1_SIZE)
    l1_circles = []
    for name, (left, top, size, cx, cy) in circle_info.items():
        if abs(size - target_size) < 1:  # Allow small tolerance
            l1_circles.append(name)

    if len(l1_circles) == 0:
        return  # No circles of target size

    # Get slide dimensions
    prs = None
    try:
        prs = slide.part.package.presentation if hasattr(slide.part, 'package') else None
    except Exception:
        prs = None

    if prs is None:
        # Fallback defaults (10in x 7.5in)
        slide_width = Inches(10)
        slide_height = Inches(7.5)
    else:
        slide_width = prs.slide_width
        slide_height = prs.slide_height

    # Calculate horizontal spacing with equal white space at boundaries
    num_circles = len(l1_circles)
    
    # Total width occupied by circles
    total_circle_width = num_circles * target_size
    
    # Remaining space to distribute (left margin + gaps between circles + right margin)
    remaining_space = slide_width - total_circle_width
    
    # Equal spacing: left boundary, between circles, and right boundary
    # Total gaps = num_circles + 1 (1 left + (num_circles - 1) between + 1 right)
    gap_size = remaining_space / (num_circles + 1)
    
    # Starting position (first gap from left)
    start_x = gap_size
    center_y = (slide_height - target_size) / 2

    # Reposition each circle
    for idx, name in enumerate(l1_circles):
        new_left = start_x + idx * (target_size + gap_size)
        new_top = center_y
        new_cx = new_left + target_size / 2
        new_cy = new_top + target_size / 2

        # Update circle_info
        circle_info[name] = (new_left, new_top, target_size, new_cx, new_cy)

        # Find and update the shape
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape.text_frame.text == name:
                shape.left = int(new_left)
                shape.top = int(new_top)
                break


def center_tool_nodes(slide) -> None:
    """Find all circles of size 0.5 inches and arrange them in a row at top of slide.

    This function identifies circles with size 0.5 inches in circle_info and
    repositions them horizontally centered in the top portion of the slide
    with equal spacing at boundaries.

    Args:
        slide: The slide object with circle_info attribute.
    """
    if not hasattr(slide, 'circle_info') or not isinstance(slide.circle_info, dict):
        print("Error: No circle information found on slide.")
        return

    circle_info = slide.circle_info

    # Find all circles with size 0.5 inches
    target_size = Inches(TOOL_SIZE)
    tool_circles = []
    for name, (left, top, size, cx, cy) in circle_info.items():
        if abs(size - target_size) < 1:  # Allow small tolerance
            tool_circles.append(name)

    if len(tool_circles) == 0:
        return  # No circles of target size

    # Get slide dimensions
    prs = None
    try:
        prs = slide.part.package.presentation if hasattr(slide.part, 'package') else None
    except Exception:
        prs = None

    if prs is None:
        # Fallback defaults (10in x 7.5in)
        slide_width = Inches(10)
        slide_height = Inches(7.5)
    else:
        slide_width = prs.slide_width
        slide_height = prs.slide_height

    # Calculate horizontal spacing with equal white space at boundaries
    num_circles = len(tool_circles)
    
    # Total width occupied by circles
    total_circle_width = num_circles * target_size
    
    # Remaining space to distribute
    remaining_space = slide_width - total_circle_width
    
    # Equal spacing: left boundary, between circles, and right boundary
    gap_size = remaining_space / (num_circles + 1)
    
    # Starting position (first gap from left) and top portion positioning
    start_x = gap_size
    top_y = Inches(1.0)  # Position in top part of slide

    # Reposition each circle
    for idx, name in enumerate(tool_circles):
        new_left = start_x + idx * (target_size + gap_size)
        new_top = top_y
        new_cx = new_left + target_size / 2
        new_cy = new_top + target_size / 2

        # Update circle_info
        circle_info[name] = (new_left, new_top, target_size, new_cx, new_cy)

        # Find and update the shape
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape.text_frame.text == name:
                shape.left = int(new_left)
                shape.top = int(new_top)
                break


def write_surround_circles(slide, source: str, dest: list) -> None:
    """Position or create circles around a source circle based on dest list.

    For each item name in dest, if a circle with that name exists, move it to
    1.5 inches from the source center. If it doesn't exist, create a new circle
    of size AGENT_SIZE at that position.

    Args:
        slide: The slide object with circle_info attribute.
        source: The name of the source circle to surround.
        dest: A list of circle names to position around the source.
    """
    if not hasattr(slide, 'circle_info') or not isinstance(slide.circle_info, dict):
        print("Error: No circle information found on slide.")
        return

    circle_info = slide.circle_info

    if source not in circle_info:
        print(f"Error: Circle '{source}' not found on slide.")
        return

    # Get source circle info
    left_src, top_src, size_src, cx_src, cy_src = circle_info[source]

    print ('cx_src ' + str(cx_src) + ' cy_src ' + str(cy_src))

    radius_distance = Inches(1.5)  # 1.5 inches from source center

    # Deduplicate dest list while preserving order
    seen = set()
    unique_dest = []
    for name in dest:
        if name not in seen:
            seen.add(name)
            unique_dest.append(name)

    if len(unique_dest) == 0:
        return  # No circles to process

    # Filter to only include circles that already exist
    existing_circles = [name for name in unique_dest if name in circle_info]

    if len(existing_circles) == 0:
        return  # No existing circles to move

    # Process only existing circles to position them
    num_circles = len(existing_circles)
    
    for i, dest_name in enumerate(existing_circles):
        angle = (2 * math.pi * i) / num_circles
        # Calculate position at 1.5 inches from source center
        x_offset = radius_distance * math.cos(angle)
        y_offset = radius_distance * math.sin(angle)

        # Calculate center position for the circle
        new_cx = cx_src + x_offset
        new_cy = cy_src + y_offset
        print ('new_cx ' + str(new_cx) + ' new_cy ' + str(new_cy))

        # Circle exists - move it
        _, _, existing_size, _, _ = circle_info[dest_name]
        new_left = new_cx - existing_size / 2
        new_top = new_cy - existing_size / 2

        # Update circle_info
        circle_info[dest_name] = (new_left, new_top, existing_size, new_cx, new_cy)

        # Find and move the shape
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape.text_frame.text == dest_name:
                shape.left = int(new_left)
                shape.top = int(new_top)
                break

